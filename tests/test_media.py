"""Tests for the media extraction service — docs/spec-media-extraction.md.

Covers what runs without LibreOffice: sanitisation (§5), clustering (§3.2),
unit conversion (§3.2, §8.2) and PPTX enumeration (§3.1). The render path needs
LibreOffice on the host and is exercised in the image, not here.
"""
import io

import pytest
from pptx import Presentation
from pptx.util import Emu, Inches

from app_media.cluster import Box, cluster_boxes, filter_clusters
from app_media.derivatives import crop_src_rect, normalise
from app_media.pptx_scan import scan_pptx
from app_media.sanitise import SanitisationError, sanitise_svg
from app_media.units import RectEmu, emu_to_pt, resolution_adequacy

SVG = '<svg xmlns="http://www.w3.org/2000/svg" width="10" height="10">{}</svg>'


# ── §5 sanitisation ─────────────────────────────────────────────────────────

def test_script_and_onload_are_removed_and_reported():
    """§11.4 — the acceptance criterion, verbatim."""
    raw = SVG.format('<script>alert(1)</script><rect onload="alert(2)" width="5" height="5"/>')
    clean, report = sanitise_svg(raw.encode(), "a1")
    text = clean.decode()
    assert "script" not in text
    assert "alert(1)" not in text          # the body went with the element
    assert "onload" not in text
    assert "script" in report["removed"]
    assert any(r.endswith("onload") for r in report["removed"])


def test_the_picture_survives_sanitisation():
    """§11.12 — criterion 4 proves the script is gone, not that anything is left."""
    raw = SVG.format('<script>x</script><path d="M0 0 L10 10"/><text>Fläche</text>')
    clean, report = sanitise_svg(raw.encode(), "a1")
    assert b"<path" in clean and "Fläche" in clean.decode()
    assert report["drawingElementCount"] == 2


def test_foreign_object_is_removed_with_its_children():
    raw = SVG.format('<foreignObject><div onclick="x">hi</div></foreignObject><rect/>')
    clean, _ = sanitise_svg(raw.encode(), "a1")
    assert b"foreignObject" not in clean and b"onclick" not in clean and b"hi" not in clean


def test_animate_elements_are_removed():
    raw = SVG.format('<rect><animate attributeName="x"/><animateTransform/></rect>')
    clean, report = sanitise_svg(raw.encode(), "a1")
    assert b"animate" not in clean
    assert {"animate", "animateTransform"} <= set(report["removed"])


@pytest.mark.parametrize("href,kept", [
    ("#frag", True),
    ("data:image/png;base64,AAAA", True),
    ("data:image/jpeg;base64,AAAA", True),
    ("javascript:alert(1)", False),
    ("https://evil.example/x.svg", False),
    ("data:text/html,alert(1)", False),
])
def test_href_allowlist(href, kept):
    raw = SVG.format(f'<image href="{href}"/>')
    clean, _ = sanitise_svg(raw.encode(), "a1")
    assert (href in clean.decode()) is kept


def test_css_import_and_off_document_url_are_stripped():
    raw = SVG.format('<style>@import url(//evil);.a{fill:url(https://evil/x)}</style>'
                     '<rect style="fill:url(https://evil/y)"/>')
    clean, _ = sanitise_svg(raw.encode(), "a1")
    assert b"@import" not in clean and b"evil" not in clean


def test_ids_are_namespaced_and_references_follow():
    raw = SVG.format('<defs><linearGradient id="g"/></defs><rect fill="url(#g)"/>'
                     '<use href="#g"/>')
    clean, _ = sanitise_svg(raw.encode(), "ast7")
    text = clean.decode()
    assert 'id="ast7-g"' in text and "url(#ast7-g)" in text and 'href="#ast7-g"' in text


def test_unknown_element_is_unwrapped_not_dropped():
    raw = SVG.format('<madeUp><rect width="1" height="1"/></madeUp>')
    clean, report = sanitise_svg(raw.encode(), "a1")
    assert b"madeUp" not in clean and b"<rect" in clean
    assert "madeUp" in report["removed"]


def test_cyclic_use_is_rejected():
    raw = SVG.format('<symbol id="a"><use href="#b"/></symbol>'
                     '<symbol id="b"><use href="#a"/></symbol>')
    with pytest.raises(SanitisationError, match="cyclic"):
        sanitise_svg(raw.encode(), "a1")


def test_non_svg_root_is_rejected():
    with pytest.raises(SanitisationError, match="expected <svg>"):
        sanitise_svg(b"<html><body/></html>", "a1")


def test_malformed_xml_is_rejected():
    with pytest.raises(SanitisationError, match="not parseable"):
        sanitise_svg(b"<svg><rect>", "a1")


# ── §3.2 clustering ─────────────────────────────────────────────────────────

def _box(key, l, t, w, h, **kw):
    return Box(key=key, rect=RectEmu(l, t, w, h), has_fill=kw.pop("fill", True), **kw)


IN = 914400
PAGE = RectEmu(0, 0, 26 * IN, 15 * IN)


def test_adjacent_shapes_merge_into_one_cluster():
    boxes = [_box("1", 0, 0, IN, IN), _box("2", IN + IN // 8, 0, IN, IN)]  # 0.125 in apart
    assert len(cluster_boxes(boxes)) == 1


def test_distant_shapes_stay_separate():
    boxes = [_box("1", 0, 0, IN, IN), _box("2", 5 * IN, 0, IN, IN)]
    assert len(cluster_boxes(boxes)) == 2


def test_a_bridging_shape_merges_two_clusters():
    """The fixed-point pass: without it, order decides the outcome."""
    left = _box("1", 0, 0, IN, IN)
    right = _box("3", 3 * IN, 0, IN, IN)
    bridge = _box("2", int(1.1 * IN), 0, int(1.85 * IN), IN)
    assert len(cluster_boxes([left, right, bridge])) == 1


def test_clustering_is_order_independent():
    """§2 promises idempotency on the source hash; §3.2 step 2 is order-dependent
    unless the input is sorted. This is that guarantee."""
    boxes = [_box(str(i), i * IN // 2, 0, IN // 3, IN) for i in range(6)]
    a = [sorted(c.keys) for c in cluster_boxes(boxes)]
    b = [sorted(c.keys) for c in cluster_boxes(list(reversed(boxes)))]
    assert a == b


def test_pure_text_cluster_is_discarded():
    """§3.2 step 4 — a bullet layout is not a diagram."""
    boxes = [_box(str(i), 0, i * IN, 4 * IN, IN // 2, fill=False, is_text=True) for i in range(4)]
    assert filter_clusters(cluster_boxes(boxes), PAGE) == []


def test_small_cluster_is_discarded():
    boxes = [_box(str(i), i * IN // 20, 0, IN // 20, IN // 20) for i in range(4)]
    assert filter_clusters(cluster_boxes(boxes), PAGE) == []


def test_three_filled_shapes_of_sufficient_area_survive():
    boxes = [_box(str(i), i * 3 * IN, 0, int(2.8 * IN), 6 * IN) for i in range(3)]
    kept = filter_clusters(cluster_boxes(boxes), PAGE)
    assert len(kept) == 1 and len(kept[0].boxes) == 3


# ── §3.2 / §8.2 units ───────────────────────────────────────────────────────

def test_emu_to_point_is_exact():
    assert emu_to_pt(12700) == 1.0
    assert emu_to_pt(round(26.67 * 914400)) == pytest.approx(1920, abs=0.5)


def test_gap_is_zero_when_boxes_overlap():
    assert RectEmu(0, 0, 100, 100).gap_to(RectEmu(50, 50, 100, 100)) == 0


def test_resolution_adequacy_accounts_for_css_pixel_density():
    """The 96/72 factor is not optional: without it the formula understates the
    requirement by 33% and passes images that are only ~1.5x dense (§8.2)."""
    # 400 pt wide at 2x density needs 400 * 96/72 * 2 = 1066 px
    assert resolution_adequacy(1066, 400) == pytest.approx(1.0, abs=0.01)
    # 700 px at 400 pt: the draft formula (no 96/72) reports 0.875 and passes;
    # with the conversion it is 0.66 and is correctly flagged as low resolution.
    assert 700 / (400 * 2) > 0.75
    assert resolution_adequacy(700, 400) < 0.75
    assert resolution_adequacy(100, 0) == 0.0            # no division by zero


# ── §3.1 PPTX enumeration ───────────────────────────────────────────────────

def _deck_with(build) -> str:
    import tempfile

    prs = Presentation()
    prs.slide_width, prs.slide_height = Emu(26 * IN), Emu(15 * IN)
    build(prs)
    path = tempfile.mktemp(suffix=".pptx")
    prs.save(path)
    return path


def test_diagram_like_slide_yields_a_shape_group():
    from pptx.enum.shapes import MSO_SHAPE

    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        for i in range(3):
            slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1 + i * 3), Inches(3), Inches(2.9), Inches(6))

    scans, page, _ = scan_pptx(_deck_with(build))
    groups = [c for c in scans[0].candidates if c.cls == "shape_group"]
    assert len(groups) == 1
    assert len(groups[0].shape_ids) == 3


def test_theme_filled_shape_counts_as_filled():
    """Regression: python-pptx reports `fill.type is None` for a shape whose fill
    comes from the theme via p:style/a:fillRef — the normal case for a drawn
    diagram box. Trusting it made §3.2 step 4 discard exactly the shapes this
    pipeline exists to find."""
    from pptx.enum.shapes import MSO_SHAPE

    from app_media.pptx_scan import _has_fill

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                 Inches(1), Inches(1), Inches(2), Inches(2))
    text = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(1))

    assert box.fill.type is None          # what misled the first implementation
    assert _has_fill(box) is True         # fillRef → drawn
    assert _has_fill(text) is False       # explicit a:noFill → not drawn


def test_bullet_slide_yields_no_candidate():
    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Nur Text"
        slide.placeholders[1].text_frame.text = "a\nb\nc"

    scans, _, _ = scan_pptx(_deck_with(build))
    assert scans[0].candidates == []


def test_context_is_extracted():
    """§7.1 — context is the more valuable half, and it is extracted, not inferred."""
    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Die vier Werkzeuge"
        slide.placeholders[1].text_frame.text = "Erstens ..."
        slide.notes_slide.notes_text_frame.text = "Hier die Analogie erklären"

    scans, _, _ = scan_pptx(_deck_with(build))
    assert scans[0].title == "Die vier Werkzeuge"
    assert "Erstens" in scans[0].surrounding_text
    assert "Analogie" in scans[0].speaker_notes


# ── §3.1 srcRect ────────────────────────────────────────────────────────────

def test_src_rect_crop_matches_the_on_slide_crop():
    """§11.7 — registering the uncrop produces assets nobody recognises later."""
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (1000, 500), "white").save(buf, format="PNG")
    image = normalise(buf.getvalue())
    cropped = crop_src_rect(image, {"l": 10000, "r": 10000, "t": 0, "b": 50000})
    assert cropped.size == (800, 250)


def test_no_src_rect_leaves_the_image_alone():
    from PIL import Image

    buf = io.BytesIO()
    Image.new("RGB", (100, 100)).save(buf, format="PNG")
    image = normalise(buf.getvalue())
    assert crop_src_rect(image, None).size == (100, 100)


# ── §4.3 SVG crop ───────────────────────────────────────────────────────────

PAGE_SVG = (
    '<svg xmlns="http://www.w3.org/2000/svg" viewBox="0 0 1920 1080" '
    'width="1920pt" height="1080pt">'
    '<rect id="inside" x="100" y="100" width="200" height="200"/>'
    '<rect id="far" x="1700" y="900" width="50" height="50"/>'
    '<rect id="straddle" x="280" y="150" width="200" height="50"/>'
    '<text x="120" y="150">Fläche</text>'
    '</svg>'
).encode()


def _crop_box(l_pt, t_pt, w_pt, h_pt):
    from app_media.units import EMU_PER_POINT
    return RectEmu(l_pt * EMU_PER_POINT, t_pt * EMU_PER_POINT,
                   w_pt * EMU_PER_POINT, h_pt * EMU_PER_POINT)


def test_crop_sets_the_viewbox_without_rescaling_geometry():
    """§4.3 step 3 — the slide's points map 1:1 onto PDF user space, so the
    bounding box *is* the crop rectangle. Rescaling would break that."""
    from app_media.svg_crop import crop_svg

    out = crop_svg(PAGE_SVG, _crop_box(100, 100, 300, 200)).decode()
    assert 'viewBox="100.00 100.00 300.00 200.00"' in out
    assert 'x="100"' in out          # the inside rect kept its original coordinates


def test_crop_drops_elements_fully_outside_and_keeps_straddlers():
    """§4.3 step 4 — half a connector beats a diagram missing its arrows."""
    from app_media.svg_crop import crop_svg

    out = crop_svg(PAGE_SVG, _crop_box(100, 100, 300, 200)).decode()
    assert 'id="far"' not in out         # entirely outside
    assert 'id="straddle"' in out        # crosses the edge — kept and clipped
    assert "clip-path" in out and "clipPath" in out


def test_crop_rejects_a_degenerate_box():
    from app_media.svg_crop import CropError, crop_svg

    with pytest.raises(CropError, match="degenerate"):
        crop_svg(PAGE_SVG, _crop_box(0, 0, 0, 100))


def test_cropped_output_survives_sanitisation():
    from app_media.svg_crop import crop_svg

    clean, report = sanitise_svg(crop_svg(PAGE_SVG, _crop_box(100, 100, 300, 200)), "a1")
    assert report["drawingElementCount"] >= 2
    assert "Fläche" in clean.decode()


# ── §3.3 / §11.3 PDF path ───────────────────────────────────────────────────

def _one_page_pdf(path):
    import fitz

    doc = fitz.open()
    page = doc.new_page(width=1920, height=1080)
    page.insert_text((100, 200), "Liegenschaft Weidgasse", fontsize=28)
    # 10 pt apart. The GAP threshold is 0.25 in = 18 pt, so 260-spacing (20 pt
    # gaps) would *not* merge — a reminder that this constant is a guess until
    # the §11.1 hand-labelled set exists: 20 pt reads as one diagram to a human.
    for i in range(3):
        page.draw_rect(fitz.Rect(100 + i * 250, 300, 340 + i * 250, 700),
                       color=(0, 0, 0), fill=(0.9, 0.9, 0.9))
    doc.save(path)
    doc.close()
    return path


def test_page_to_svg_keeps_text_as_text(tmp_path):
    """§11.3 — the acceptance criterion, and a real trap: PyMuPDF's *default*
    converts every glyph to a path, which is exactly the failure §4.3 warns about
    for pdftocairo. Output that is not searchable is worthless to a screen reader.
    """
    from app_media.pdf_scan import page_to_svg

    svg = page_to_svg(str(_one_page_pdf(tmp_path / "t.pdf")), 1).decode()
    assert "<text" in svg
    assert "Liegenschaft" in svg


def test_pdf_enumeration_finds_a_drawing_cluster(tmp_path):
    from app_media.pdf_scan import scan_pdf

    scans, page = scan_pdf(str(_one_page_pdf(tmp_path / "t.pdf")))
    assert page.w > 0
    candidates = scans[0].candidates
    assert candidates and candidates[0].cls == "shape_group"


def test_pdf_region_render_crops(tmp_path):
    from app_media.pdf_scan import render_region

    png = render_region(str(_one_page_pdf(tmp_path / "t.pdf")), 1,
                        _crop_box(100, 300, 240, 400), dpi=72)
    image = normalise(png)
    assert image.size == pytest.approx((240, 400), abs=2)


def test_header_band_rejection_needs_repetition():
    """§3.3 — a narrow cluster in the footer band is furniture only if it repeats."""
    from app_media.pdf_scan import _in_band

    page = RectEmu(0, 0, 1000, 1000)
    assert _in_band(RectEmu(0, 950, 100, 30), page) is True     # narrow, in the band
    assert _in_band(RectEmu(0, 950, 400, 30), page) is False    # wide enough to be content
    assert _in_band(RectEmu(0, 400, 100, 30), page) is False    # mid-page


# ── §5 decode guard ─────────────────────────────────────────────────────────

def test_ordinary_high_dpi_screenshot_is_accepted():
    """Regression: the guard was set at 20 MB decoded, i.e. 5.2 MP — below a
    normal screenshot. Five figures in a real 128-slide deck were dropped by it,
    including a 3713x2475 one. It is a decompression-bomb guard, not a content
    limit."""
    from PIL import Image

    from app_media.derivatives import MAX_DECODED_PIXELS, normalise

    assert 3713 * 2475 < MAX_DECODED_PIXELS

    buf = io.BytesIO()
    Image.new("RGB", (3713, 2475), "white").save(buf, format="PNG")
    assert normalise(buf.getvalue()).size == (3713, 2475)


def test_a_decompression_bomb_is_still_refused():
    from app_media.derivatives import MAX_DECODED_PIXELS, DerivativeError, _open_bounded

    class FakeHuge:
        size = (50000, 50000)

        def load(self):
            raise AssertionError("must be refused from the header, before decoding")

    assert 50000 * 50000 > MAX_DECODED_PIXELS
    import app_media.derivatives as d

    original = d.Image.open
    d.Image.open = lambda _b: FakeHuge()
    try:
        with pytest.raises(DerivativeError, match="decode limit"):
            _open_bounded(b"x")
    finally:
        d.Image.open = original


# ── §3.2 whole-slide rule ───────────────────────────────────────────────────

def test_a_slide_with_primitives_yields_exactly_one_candidate():
    """Clustering split single diagrams into their parts on real decks: a
    diagram's boxes are often further apart than the gap threshold, so one
    illustration came out as four assets. A reader sees one picture."""
    from pptx.enum.shapes import MSO_SHAPE

    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        # Deliberately far apart — the old clustering would make three candidates.
        for i in range(3):
            slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(1 + i * 8), Inches(3), Inches(2), Inches(5))

    scans, _, _ = scan_pptx(_deck_with(build))
    groups = [c for c in scans[0].candidates if c.cls == "shape_group"]
    assert len(groups) == 1
    assert len(groups[0].shape_ids) == 3          # all of them, in one box


def test_furniture_is_excluded_from_the_candidate():
    """§3.2 step 1 — "non-placeholder shapes". Without it every figure comes out
    with the slide title and the page number attached."""
    from pptx.enum.shapes import MSO_SHAPE

    from app_media.units import emu_to_pt

    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[5])   # title only
        slide.shapes.title.text = "Überschrift"
        slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(4), Inches(6), Inches(6), Inches(5))

    scans, _, _ = scan_pptx(_deck_with(build))
    groups = [c for c in scans[0].candidates if c.cls == "shape_group"]
    assert len(groups) == 1
    # The box starts at the drawn shape, not at the title above it.
    assert emu_to_pt(groups[0].rect.t) > emu_to_pt(Inches(5))


def test_a_slide_without_primitives_yields_nothing():
    def build(prs):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Nur Text"
        slide.placeholders[1].text_frame.text = "a\nb\nc"

    scans, _, _ = scan_pptx(_deck_with(build))
    assert [c for c in scans[0].candidates if c.cls == "shape_group"] == []
