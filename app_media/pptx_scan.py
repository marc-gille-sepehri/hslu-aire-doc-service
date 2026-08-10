"""PPTX enumeration — spec-media-extraction.md §3.1 and §3.2.

Walks slides, not `ppt/media/`. The part directory has no slide attribution and no
way to tell a diagram from a footer mark; the slide XML has both, and resolving
images through the slide's rels is what makes the layout filter possible at all.
"""
from __future__ import annotations

import hashlib
from collections import defaultdict
from dataclasses import dataclass, field
from typing import Iterator

from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Emu

from .cluster import Box, cluster_boxes, filter_clusters
from .units import RectEmu

# §3.1 frequency filter
FREQUENCY_FRACTION = 0.10
FREQUENCY_FLOOR = 3

SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"   # a:extLst → svgBlip
DIAGRAM_URI = "http://schemas.openxmlformats.org/drawingml/2006/diagram"
CHART_URI = "http://schemas.openxmlformats.org/drawingml/2006/chart"


@dataclass
class Candidate:
    slide: int
    cls: str                       # svg_native | vector_import | chart | shape_group | raster
    rect: RectEmu
    shape_ids: list[str]
    author_alt_text: str | None = None
    shape_name: str | None = None
    part_name: str | None = None   # for image-backed classes
    content_hash: str | None = None
    src_rect: dict | None = None   # §3.1 — the on-slide crop, insets in 1/100000
    contained_smartart: bool = False
    native_pixels: tuple[int, int] | None = None

    def as_dict(self) -> dict:
        return {
            "slide": self.slide,
            "class": self.cls,
            "boundingBoxEmu": self.rect.as_dict(),
            "boundingBoxPt": self.rect.as_points(),
            "shapeIds": self.shape_ids,
            "authorAltText": self.author_alt_text,
            "shapeName": self.shape_name,
            "srcRect": self.src_rect,
            "containedSmartArt": self.contained_smartart,
        }


@dataclass
class SlideScan:
    number: int
    title: str | None
    surrounding_text: str
    speaker_notes: str
    layout_name: str | None
    candidates: list[Candidate] = field(default_factory=list)


def _rect(shape) -> RectEmu | None:
    if shape.left is None or shape.top is None or shape.width is None or shape.height is None:
        return None
    return RectEmu(int(shape.left), int(shape.top), int(shape.width), int(shape.height))


def _rotated_extent(shape) -> RectEmu | None:
    """§3.2: use the rotated extent, or one rotated arrow drags in half a slide."""
    base = _rect(shape)
    if base is None:
        return None
    try:
        rot = float(shape.rotation or 0.0)
    except (AttributeError, TypeError, ValueError):
        return base
    if abs(rot % 180.0) < 0.01:
        return base
    import math

    rad = math.radians(rot)
    cx, cy = base.l + base.w / 2, base.t + base.h / 2
    w = abs(base.w * math.cos(rad)) + abs(base.h * math.sin(rad))
    h = abs(base.w * math.sin(rad)) + abs(base.h * math.cos(rad))
    return RectEmu(round(cx - w / 2), round(cy - h / 2), round(w), round(h))


_EXPLICIT_FILLS = ("a:solidFill", "a:gradFill", "a:blipFill", "a:pattFill")


def _has_fill(shape) -> bool:
    """§3.2 step 4 — "no shape has a non-noFill fill".

    Read the XML, not `shape.fill.type`: python-pptx returns None for a shape
    whose fill comes from the theme via `p:style/a:fillRef`, which is the normal
    case for a drawn diagram box. Trusting it discards exactly the shapes this
    pipeline exists to find.
    """
    el = shape._element  # noqa: SLF001
    spPr = el.find(qn("p:spPr"))
    if spPr is not None and spPr.find(qn("a:noFill")) is not None:
        return False
    if spPr is not None and any(spPr.find(qn(tag)) is not None for tag in _EXPLICIT_FILLS):
        return True
    return el.find(f".//{qn('a:fillRef')}") is not None


def _alt_text(shape) -> str | None:
    """§3.1 — p:cNvPr/@descr and @title. Author-supplied, so it outranks generation."""
    try:
        nv = shape._element._nvXxPr.cNvPr  # noqa: SLF001 — python-pptx exposes no public accessor
    except Exception:
        return None
    for attr in ("descr", "title"):
        value = (nv.get(attr) or "").strip()
        if value:
            return value
    return None


def _blip_of(shape):
    blips = shape._element.findall(f".//{qn('a:blip')}")  # noqa: SLF001
    return blips[0] if blips else None


def _is_svg_native(blip) -> bool:
    """PowerPoint stores an SVG alongside a PNG fallback; the SVG wins (§3)."""
    if blip is None:
        return False
    for ext in blip.findall(f".//{qn('a:ext')}"):
        if ext.get("uri") == SVG_EXT_URI:
            return True
    return False


def _src_rect(shape) -> dict | None:
    el = shape._element.find(f".//{qn('a:srcRect')}")  # noqa: SLF001
    if el is None:
        return None
    insets = {k: int(el.get(k, 0)) for k in ("l", "t", "r", "b")}
    return insets if any(insets.values()) else None


def _graphic_uri(shape) -> str | None:
    data = shape._element.find(f".//{qn('a:graphicData')}")  # noqa: SLF001
    return data.get("uri") if data is not None else None


def _classify_picture(shape, image) -> str:
    """§3 resolution order. `svg_native` first, deliberately: PowerPoint stores a
    PNG fallback beside the SVG, and classifying the fallback first would give the
    pair two different content hashes so neither reaches the frequency threshold."""
    if _is_svg_native(_blip_of(shape)):
        return "svg_native"
    ext = (getattr(image, "ext", "") or "").lower().lstrip(".")
    if ext in ("emf", "wmf"):
        return "vector_import"
    return "raster"


def _iter_shapes(shapes) -> Iterator:
    """Yield shapes, treating a group as one unit — §3.2, only the outermost group."""
    for shape in shapes:
        yield shape


def _slide_text(slide) -> tuple[str | None, str]:
    title, parts = None, []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        if title is None and shape == getattr(slide.shapes, "title", None):
            title = text
        else:
            parts.append(text)
    return title, "\n".join(parts)[:2000]


def scan_pptx(path: str) -> tuple[list[SlideScan], RectEmu, list[dict]]:
    """Return (slides, slide_rect, rejections).

    `rejections` carries the frequency-filtered hashes with their slide lists —
    §3.1 notes that a diagram deliberately reused as a running reference is
    exactly what this filter wrongly removes, so the review pass has to see them.
    """
    prs = Presentation(path)
    page = RectEmu(0, 0, int(prs.slide_width), int(prs.slide_height))

    scans: list[SlideScan] = []
    hash_slides: dict[str, set[int]] = defaultdict(set)

    for index, slide in enumerate(prs.slides, start=1):
        title, surrounding = _slide_text(slide)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = (slide.notes_slide.notes_text_frame.text or "").strip()[:2000]

        scan = SlideScan(
            number=index,
            title=title,
            surrounding_text=surrounding,
            speaker_notes=notes,
            layout_name=getattr(slide.slide_layout, "name", None),
        )

        boxes: list[Box] = []
        for shape in _iter_shapes(slide.shapes):
            rect = _rotated_extent(shape)
            if rect is None:
                continue
            shape_id = str(shape.shape_id)
            uri = _graphic_uri(shape) if shape.shape_type is not None else None

            # Picture-backed candidates are their own class and skip clustering.
            image = getattr(shape, "image", None)
            if image is not None:
                blob = getattr(image, "blob", b"")
                digest = hashlib.sha256(blob).hexdigest() if blob else None
                if digest:
                    hash_slides[digest].add(index)
                scan.candidates.append(
                    Candidate(
                        slide=index,
                        cls=_classify_picture(shape, image),
                        rect=rect,
                        shape_ids=[shape_id],
                        author_alt_text=_alt_text(shape),
                        shape_name=shape.name,
                        content_hash=digest,
                        src_rect=_src_rect(shape),
                        native_pixels=getattr(image, "size", None),
                    )
                )
                continue

            if uri == CHART_URI:
                scan.candidates.append(
                    Candidate(index, "chart", rect, [shape_id],
                              _alt_text(shape), shape.name)
                )
                continue

            is_diagram = uri == DIAGRAM_URI
            boxes.append(
                Box(
                    key=shape_id,
                    rect=rect,
                    has_fill=_has_fill(shape),
                    is_connector=bool(getattr(shape, "connector_type", None)),
                    is_text=bool(getattr(shape, "has_text_frame", False)),
                    preformed_group=bool(getattr(shape, "shapes", None)) or is_diagram,
                )
            )

        for cluster in filter_clusters(cluster_boxes(boxes), page):
            first = cluster.boxes[0]
            scan.candidates.append(
                Candidate(
                    slide=index,
                    cls="shape_group",
                    rect=cluster.rect,
                    shape_ids=cluster.keys,
                    contained_smartart=first.preformed_group and len(cluster.boxes) == 1,
                )
            )

        scans.append(scan)

    rejections = _apply_frequency_filter(scans, hash_slides, len(scans))
    return scans, page, rejections


def extract_media(path: str, slide: int, shape_id: str) -> tuple[bytes, str]:
    """Return (bytes, extension) for a picture-backed candidate.

    For `svg_native` this resolves the `svgBlip` in the blip's `a:extLst` and
    returns the **SVG**, not the PNG fallback PowerPoint stores beside it (§3).
    Emitting the fallback would register a raster where a vector exists — the
    single most avoidable fidelity loss in the whole pipeline.
    """
    prs = Presentation(path)
    slides = list(prs.slides)
    if not 1 <= slide <= len(slides):
        raise KeyError(f"slide {slide} out of range")
    for shape in slides[slide - 1].shapes:
        if str(shape.shape_id) != str(shape_id):
            continue
        blip = _blip_of(shape)
        if blip is not None:
            svg = _svg_part(shape, blip)
            if svg is not None:
                return svg, "svg"
        image = getattr(shape, "image", None)
        if image is None:
            raise KeyError(f"shape {shape_id} on slide {slide} carries no media")
        return image.blob, (image.ext or "bin").lstrip(".").lower()
    raise KeyError(f"shape {shape_id} not found on slide {slide}")


def _svg_part(shape, blip) -> bytes | None:
    """Resolve a:extLst → asvg:svgBlip → r:embed → the SVG part's bytes."""
    for ext in blip.findall(f".//{qn('a:ext')}"):
        if ext.get("uri") != SVG_EXT_URI:
            continue
        for child in ext.iter():
            rid = child.get(qn("r:embed")) or child.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            if not rid:
                continue
            try:
                return shape.part.related_part(rid).blob
            except Exception:
                return None
    return None


def _apply_frequency_filter(scans: list[SlideScan], hash_slides: dict, slide_count: int) -> list[dict]:
    """§3.1 — drop decoration that survived the layout filter, and report what was dropped."""
    threshold = max(FREQUENCY_FLOOR, round(slide_count * FREQUENCY_FRACTION))
    rejected = {h for h, slides in hash_slides.items() if len(slides) >= threshold}
    report = []
    for digest in sorted(rejected):
        report.append({"contentHash": digest, "slides": sorted(hash_slides[digest]),
                       "threshold": threshold})
    for scan in scans:
        keep = []
        for candidate in scan.candidates:
            if candidate.content_hash in rejected and not candidate.author_alt_text:
                continue
            keep.append(candidate)
        scan.candidates = keep
    return report
